"""
Document parser: converts raw bytes → cleaned text + metadata.
Supports PDF, DOCX, PPTX, Markdown, plain text, HTML (email/Notion).
"""
from __future__ import annotations

import hashlib
import io
import re
from collections import defaultdict
from datetime import datetime
from typing import Any

from langdetect import detect, LangDetectException


# ─── Public interface ────────────────────────────────────────────────────────

class ParsedDocument:
    def __init__(
        self,
        text: str,
        title: str | None,
        author: str | None,
        language: str,
        creation_date: datetime | None,
        modified_date: datetime | None,
        mime_type: str,
        file_hash: str,
        raw_size: int,
        metadata: dict[str, Any],
    ):
        self.text = text
        self.title = title
        self.author = author
        self.language = language
        self.creation_date = creation_date
        self.modified_date = modified_date
        self.mime_type = mime_type
        self.file_hash = file_hash
        self.raw_size = raw_size
        self.metadata = metadata


def parse_document(content: bytes, mime_type: str, filename: str = "") -> ParsedDocument:
    """Entry point: dispatch to the right parser by MIME type."""
    file_hash = hashlib.sha256(content).hexdigest()
    raw_size = len(content)

    if mime_type == "application/pdf":
        result = _parse_pdf(content)
    elif mime_type in (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    ):
        result = _parse_docx(content)
    elif mime_type in (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    ):
        result = _parse_pptx(content)
    elif mime_type in ("text/html", "application/xhtml+xml"):
        result = _parse_html(content)
    elif mime_type in ("text/markdown", "text/x-markdown"):
        result = _parse_markdown(content)
    else:
        result = {"text": content.decode("utf-8", errors="replace"), "metadata": {}}

    text = _clean_text(result["text"])
    language = _detect_language(text)
    title = result.get("title") or _infer_title(text, filename)

    return ParsedDocument(
        text=text,
        title=title,
        author=result.get("author"),
        language=language,
        creation_date=result.get("creation_date"),
        modified_date=result.get("modified_date"),
        mime_type=mime_type,
        file_hash=file_hash,
        raw_size=raw_size,
        metadata=result.get("metadata", {}),
    )


# ─── PDF ─────────────────────────────────────────────────────────────────────

def _parse_pdf(content: bytes) -> dict:
    import pdfplumber

    pages_text: list[str] = []
    all_headings: list[str] = []
    metadata: dict = {}

    with pdfplumber.open(io.BytesIO(content)) as pdf:
        if pdf.metadata:
            metadata = dict(pdf.metadata)

        for page in pdf.pages:
            text = page.extract_text(x_tolerance=2, y_tolerance=2) or ""

            # Extract tables as markdown-style rows
            for table in page.extract_tables():
                if not table:
                    continue
                rows = [" | ".join(str(c) if c else "" for c in row) for row in table]
                text += "\n\n" + "\n".join(rows)

            # Detect large-font lines and inject them as ## headings
            text = _inject_pdf_headings(page, text)
            pages_text.append(text)

            # Collect heading strings for metadata
            for line in text.split("\n"):
                if line.startswith("## "):
                    all_headings.append(line[3:].strip())

    return {
        "text": "\n\n".join(pages_text),
        "title": metadata.get("Title") or metadata.get("title"),
        "author": metadata.get("Author") or metadata.get("author"),
        "creation_date": _parse_pdf_date(metadata.get("CreationDate")),
        "modified_date": _parse_pdf_date(metadata.get("ModDate")),
        "metadata": {"headings": list(dict.fromkeys(all_headings))},
    }


def _inject_pdf_headings(page, extracted_text: str) -> str:
    """
    Scan the page's character stream to find large-font lines, then prefix
    those lines with '## ' in the already-extracted text so the chunker
    can split on them as Markdown headings.
    """
    chars = page.chars
    if not chars:
        return extracted_text

    # Group chars by line: bucket by rounded top coordinate (3-pt tolerance)
    line_buckets: dict[int, list] = defaultdict(list)
    for ch in chars:
        if ch.get("text", "").strip():
            key = round(ch.get("top", 0) / 3) * 3
            line_buckets[key].append(ch)

    # Identify heading lines: average font size > 13 pt, reasonable text length
    heading_strings: set[str] = set()
    for bucket in line_buckets.values():
        sizes = [c.get("size", 0) for c in bucket if c.get("size", 0) > 0]
        if not sizes:
            continue
        avg_size = sum(sizes) / len(sizes)
        if avg_size > 13:
            # Sort by x-coordinate to reconstruct left-to-right reading order
            text = "".join(
                c["text"] for c in sorted(bucket, key=lambda c: c.get("x0", 0))
            ).strip()
            if 2 < len(text) < 200:
                heading_strings.add(text)

    if not heading_strings:
        return extracted_text

    # Prefix matching lines in the extracted text with ##
    result_lines = []
    for line in extracted_text.split("\n"):
        stripped = line.strip()
        if stripped in heading_strings and not stripped.startswith("#"):
            result_lines.append(f"## {stripped}")
        else:
            result_lines.append(line)
    return "\n".join(result_lines)


def _parse_pdf_date(raw: str | None) -> datetime | None:
    if not raw:
        return None
    try:
        raw = raw.replace("D:", "").replace("'", "").replace("Z", "+00:00")
        return datetime.strptime(raw[:14], "%Y%m%d%H%M%S")
    except Exception:
        return None


# ─── DOCX ────────────────────────────────────────────────────────────────────

def _parse_docx(content: bytes) -> dict:
    from docx import Document

    doc = Document(io.BytesIO(content))
    parts: list[str] = []
    headings: list[str] = []

    for para in doc.paragraphs:
        if not para.text.strip():
            continue

        style = para.style.name

        if style.startswith("Heading"):
            # Preserve heading level (Heading 1 → #, Heading 2 → ##, etc.)
            level_match = re.search(r"\d+", style)
            level = int(level_match.group()) if level_match else 2
            hashes = "#" * min(level, 6)
            headings.append(para.text)
            parts.append(f"\n{hashes} {para.text}\n")

        elif "List Number" in style:
            # Detect nesting level from style name (e.g. "List Number 2")
            level_match = re.search(r"\d+", style)
            depth = int(level_match.group()) - 1 if level_match else 0
            indent = "  " * depth
            parts.append(f"{indent}1. {para.text}")

        elif "List" in style:
            level_match = re.search(r"\d+", style)
            depth = int(level_match.group()) - 1 if level_match else 0
            indent = "  " * depth
            parts.append(f"{indent}- {para.text}")

        else:
            parts.append(para.text)

    # Tables
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text.strip() for cell in row.cells))

    core = doc.core_properties
    return {
        "text": "\n".join(parts),
        "title": core.title,
        "author": core.author,
        "creation_date": core.created,
        "modified_date": core.modified,
        "metadata": {"headings": headings},
    }


# ─── PPTX ────────────────────────────────────────────────────────────────────

def _parse_pptx(content: bytes) -> dict:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    slides_text: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        slide_parts = [f"## Slide {i}"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        slide_parts.append(text)
        slides_text.append("\n".join(slide_parts))

    core = prs.core_properties
    return {
        "text": "\n\n".join(slides_text),
        "title": core.title,
        "author": core.author,
        "creation_date": core.created,
        "modified_date": core.modified,
        "metadata": {"slide_count": len(prs.slides)},
    }


# ─── HTML / Email ────────────────────────────────────────────────────────────

def _parse_html(content: bytes) -> dict:
    from bs4 import BeautifulSoup
    import markdownify

    soup = BeautifulSoup(content, "lxml")

    for tag in soup(["script", "style", "nav", "footer", "aside"]):
        tag.decompose()

    title_tag = soup.find("title")
    title = title_tag.get_text(strip=True) if title_tag else None

    md = markdownify.markdownify(str(soup.body or soup), heading_style="ATX")
    return {"text": md, "title": title, "metadata": {}}


# ─── Markdown ────────────────────────────────────────────────────────────────

def _parse_markdown(content: bytes) -> dict:
    text = content.decode("utf-8", errors="replace")
    first_h1 = re.search(r"^#\s+(.+)$", text, re.MULTILINE)
    return {
        "text": text,
        "title": first_h1.group(1).strip() if first_h1 else None,
        "metadata": {},
    }


# ─── Helpers ─────────────────────────────────────────────────────────────────

def _clean_text(text: str) -> str:
    text = re.sub(r"\r\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    return text.strip()


def _detect_language(text: str) -> str:
    try:
        return detect(text[:1000])
    except LangDetectException:
        return "en"


def _infer_title(text: str, filename: str) -> str | None:
    first_line = text.strip().split("\n")[0].strip()
    if first_line and len(first_line) < 200:
        return first_line
    if filename:
        return filename.rsplit(".", 1)[0].replace("_", " ").replace("-", " ").title()
    return None
